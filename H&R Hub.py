import os
import re
from dataclasses import dataclass
from typing import List, Dict, Optional, Tuple

import streamlit as st
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openai import OpenAI

# --- CGIAR Theme --------------------------------------------------------------
import html  # to escape text in chips

import pandas as pd
import pickle
import hashlib
from dataclasses import asdict

from datetime import datetime
import json
from urllib.parse import urlparse

CGIAR_COLORS = {
    "green_primary": "#427730",      # Corporate Green
    "green_leaf": "#7AB800",         # Leaf green
    "green_leaf_dark": "#739600",    # Darker leaf green
    "blue_bright": "#0065BD",        # Bright Blue
    "blue_medium": "#0039A6",        # Medium Blue
    "yellow": "#FDC82F",             # Yellow
    "orange": "#E37222",             # Orange (amber-ish)
    "bg": "#F7FAF8",                 # Light soft background
    "panel": "#FFFFFF",              # Cards
    "text": "#1A202C",               # Main text
    "muted": "#4A5568",              # Secondary text
    "border": "#E2E8F0",             # Subtle borders
}

def apply_cgiar_theme():
    st.markdown(f"""
    <style>
        /* Typography */
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        :root {{
            --brand-primary: {CGIAR_COLORS["green_primary"]};
            --brand-primary-strong: {CGIAR_COLORS["green_leaf_dark"]};
            --brand-accent: {CGIAR_COLORS["green_leaf"]};
            --brand-blue: {CGIAR_COLORS["blue_bright"]};
            --brand-blue-strong: {CGIAR_COLORS["blue_medium"]};
            --brand-yellow: {CGIAR_COLORS["yellow"]};
            --brand-orange: {CGIAR_COLORS["orange"]};

            --bg: {CGIAR_COLORS["bg"]};
            --panel: {CGIAR_COLORS["panel"]};
            --text: {CGIAR_COLORS["text"]};
            --muted: {CGIAR_COLORS["muted"]};
            --border: {CGIAR_COLORS["border"]};
        }}

        .stApp {{
            font-family: 'Inter', system-ui, -apple-system, Segoe UI, Roboto, sans-serif;
            background: var(--bg);
            color: var(--text);
        }}
        #MainMenu {{ display: none; }}
        footer {{ visibility: hidden; }}

        /* Main container */
        .main .block-container {{
            max-width: 1100px;
            padding-top: 1.25rem;
        }}

        /* Hero — solid amber */
        .brand-hero {{
            background: var(--brand-orange);
            color: white;
            border-radius: 12px;
            padding: 1.25rem 1.25rem;
            margin-bottom: 1rem;
            box-shadow: 0 2px 8px rgba(0,0,0,0.06);
        }}
        .brand-hero h1 {{
            margin: 0 0 .25rem 0;
            font-weight: 700;
            letter-spacing: .2px;
        }}
        .brand-hero p {{
            margin: 0;
            opacity: .95;
        }}

        /* Cards */
        .card {{
            background: var(--panel);
            border: 1px solid var(--border);
            border-radius: 12px;
            padding: 1rem 1.25rem;
            box-shadow: 0 2px 8px rgba(0,0,0,0.04);
            margin: .75rem 0;
        }}
        .answer-card {{
            border-left: 4px solid var(--brand-primary);
        }}

        /* Source chips */
        .sources-wrap {{
            display: flex;
            flex-wrap: wrap;
            gap: .5rem;
            margin-top: .5rem;
        }}
        .source-chip {{
            background: #f2f7f3;
            border: 1px solid #e1efe4;
            color: #0f3b1f;
            border-radius: 999px;
            padding: .35rem .75rem;
            font-size: .85rem;
            line-height: 1;
            white-space: nowrap;
        }}

        /* Buttons */
        .stButton > button {{
            width: 100%;
            border: 0;
            border-radius: 10px;
            font-weight: 600;
            padding: .65rem 1rem;
            transition: transform .08s ease, opacity .15s ease, box-shadow .2s ease;
            background: var(--brand-primary);
            color: #fff;
            box-shadow: 0 2px 6px rgba(66,119,48,.20);
        }}
        .stButton > button:hover {{
            opacity: .95;
            box-shadow: 0 4px 10px rgba(66,119,48,.25);
            transform: translateY(-1px);
        }}
        .stButton > button:active {{
            transform: translateY(1px);
            box-shadow: inset 0 2px 4px rgba(0,0,0,.08);
        }}

        /* Inputs */
        .stTextInput > div > div > input {{
            border-radius: 10px !important;
            border: 1px solid var(--border) !important;
            box-shadow: none !important;
        }}
        .stTextInput > div > div > input:focus {{
            border-color: var(--brand-primary) !important;
            outline: 3px solid rgba(66,119,48,.15) !important;
        }}

        /* Slider */
        .stSlider [data-baseweb="slider"] > div:first-child {{
            color: var(--brand-primary);
        }}

        /* Subtle alerts */
        .stAlert {{
            border-left: 4px solid var(--brand-accent);
        }}

        /* Compact metrics */
        .metric-row > div > div {{
            background: var(--panel);
            border: 1px solid var(--border);
            border-radius: 12px;
            padding: .75rem;
        }}

        /* Footer */
        .app-footer {{
            text-align: center;
            color: var(--muted);
            font-size: .9rem;
            margin: 1rem 0 2rem;
        }}
    </style>
    """, unsafe_allow_html=True)

@dataclass
class Chunk:
    text: str
    source_path: str
    source_name: str
    kind: str  # pdf | docx | pptx
    location: str  # e.g., "page 3", "slide 2", "section X / paragraph 12"
    id: str = ""


def normalize_whitespace(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def split_sentences(text: str) -> List[str]:
    if not text:
        return []
    parts = re.split(r"(?<=[\.\!\?])\s+", text)
    sentences: List[str] = []
    for p in parts:
        p = p.strip()
        if len(p) > 0:
            sentences.append(p)
    if len(sentences) <= 1:
        lines = [ln.strip() for ln in re.split(r"[\n\r]+", text) if ln.strip()]
        if len(lines) > len(sentences):
            sentences = lines
    return sentences


def split_text_into_overlapping_parts(text: str, num_parts: int = 4, overlap_fraction: float = 0.25) -> List[str]:
    words = text.split()
    total_words = len(words)
    if total_words == 0:
        return []
    denominator = num_parts - (num_parts - 1) * overlap_fraction if num_parts > 1 else 1
    part_words = int(total_words / denominator) if denominator > 0 else total_words
    overlap_words = int(part_words * overlap_fraction)
    step = part_words - overlap_words
    parts = []
    start = 0
    for i in range(num_parts):
        end = min(start + part_words, total_words)
        part = ' '.join(words[start:end])
        if part:
            parts.append(part)
        if end >= total_words:
            break
        start += step
    return parts


def read_pdf_chunks(path: str) -> List[Chunk]:
    chunks: List[Chunk] = []
    try:
        reader = PdfReader(path)
        for index, page in enumerate(reader.pages, start=1):
            page_text = normalize_whitespace(page.extract_text() or "")
            if not page_text:
                continue
            parts = split_text_into_overlapping_parts(page_text, 4, 0.25)
            for part_idx, part in enumerate(parts, start=1):
                chunks.append(
                    Chunk(
                        text=part,
                        source_path=path,
                        source_name=os.path.basename(path),
                        kind="pdf",
                        location=f"page {index} part {part_idx}",
                        id=f"pdf:{os.path.basename(path)}:p{index}pt{part_idx}",
                    )
                )
        return chunks
    except Exception:
        return []


def read_docx_chunks(path: str) -> List[Chunk]:
    try:
        doc = DocxDocument(path)
        unit_texts: List[str] = []
        unit_locations: List[str] = []
        current_section: Optional[str] = None
        for paragraph_index, paragraph in enumerate(doc.paragraphs, start=1):
            text = normalize_whitespace(paragraph.text)
            if not text:
                continue
            style_name = getattr(paragraph.style, "name", "") or ""
            if style_name.lower().startswith("heading") or style_name.lower().startswith("título"):
                current_section = text
            location = f"section '{current_section}' paragraph {paragraph_index}" if current_section else f"paragraph {paragraph_index}"
            unit_texts.append(text)
            unit_locations.append(location)
        return create_overlapped_chunks(unit_texts, unit_locations, os.path.basename(path), "docx", path)
    except Exception:
        pass
    return []


def read_pptx_chunks(path: str) -> List[Chunk]:
    chunks: List[Chunk] = []
    try:
        prs = Presentation(path)
        for slide_index, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        txt = "\n".join(p.text for p in shape.text_frame.paragraphs)
                        txt = normalize_whitespace(txt)
                        if txt:
                            texts.append(txt)
                except Exception:
                    continue
            slide_text = normalize_whitespace("\n".join(texts))
            if not slide_text:
                continue
            parts = split_text_into_overlapping_parts(slide_text, 4, 0.25)
            for part_idx, part in enumerate(parts, start=1):
                chunks.append(
                    Chunk(
                        text=part,
                        source_path=path,
                        source_name=os.path.basename(path),
                        kind="pptx",
                        location=f"slide {slide_index} part {part_idx}",
                        id=f"pptx:{os.path.basename(path)}:s{slide_index}pt{part_idx}",
                    )
                )
        return chunks
    except Exception:
        return []


def create_overlapped_chunks(unit_texts: List[str], unit_locations: List[str], basename: str, kind: str, path: str, window_size: int = 4, overlap: int = 1) -> List[Chunk]:
    chunks: List[Chunk] = []
    step = window_size - overlap
    for i in range(0, len(unit_texts), step):
        end = i + window_size
        slice_texts = unit_texts[i:end]
        slice_locations = unit_locations[i:end] if end <= len(unit_locations) else unit_locations[i:]
        if len(slice_texts) == 0:
            continue
        merge_len = len(slice_texts)
        if merge_len < 3 and chunks:
            # merge to last
            last = chunks[-1]
            last.text += ' ' + ' '.join(slice_texts)
            last_loc_end = slice_locations[-1] if slice_locations else ""
            last.location = last.location.rsplit(" to ", 1)[0] + f" to {last_loc_end}"
            # update id
            parts = last.id.split(':')
            if len(parts) == 3 and parts[2].startswith('u'):
                range_part = parts[2][1:]
                start_str, end_str = range_part.split('-')
                start = int(start_str)
                old_end = int(end_str)
                new_end = old_end + merge_len
                last.id = f"{parts[0]}:{parts[1]}:u{start}-{new_end}"
            continue
        chunk_text = ' '.join(slice_texts)
        loc_start = slice_locations[0]
        loc_end = slice_locations[-1]
        chunk_loc = f"{loc_start} to {loc_end}"
        start_idx = i + 1
        end_idx = i + merge_len
        chunk_id = f"{kind}:{basename}:u{start_idx}-{end_idx}"
        chunks.append(
            Chunk(
                text=chunk_text,
                source_path=path,
                source_name=basename,
                kind=kind,
                location=chunk_loc,
                id=chunk_id,
            )
        )
    return chunks


def load_corpus(root_dir: str,
                hide_failures: bool = True,
                exclude_patterns: Optional[List[str]] = None) -> List[Chunk]:
    """
    Recorre root_dir y subcarpetas, carga solo .pdf/.docx/.pptx.
    Muestra únicamente los archivos cargados con éxito (por defecto).
    Puedes excluir archivos por patrón (p.ej., ['^default\\.', '^\\._', '^~\\$']).
    """
    supported_ext = {".pdf", ".docx", ".pptx"}
    exclude_patterns = exclude_patterns or []
    compiled_excludes = [re.compile(pat, re.IGNORECASE) for pat in exclude_patterns]

    chunks: List[Chunk] = []
    ok_files: List[Tuple[str, int]] = []   # (nombre, n_chunks)
    bad_files: List[str] = []

    for dirpath, _, filenames in os.walk(root_dir):
        for fname in filenames:
            # Excluir por patrón (ej.: default.*, archivos ocultos de macOS, backups temporales)
            if any(p.match(fname) for p in compiled_excludes):
                continue

            ext = os.path.splitext(fname)[1].lower()
            if ext not in supported_ext:
                continue

            abspath = os.path.join(dirpath, fname)
            # Excluir archivos vacíos
            try:
                if os.path.getsize(abspath) == 0:
                    bad_files.append(fname)
                    continue
            except Exception:
                bad_files.append(fname)
                continue

            new_chunks: List[Chunk] = []
            try:
                if ext == ".pdf":
                    new_chunks = read_pdf_chunks(abspath)
                elif ext == ".docx":
                    new_chunks = read_docx_chunks(abspath)
                elif ext == ".pptx":
                    new_chunks = read_pptx_chunks(abspath)
            except Exception:
                new_chunks = []

            if new_chunks:
                chunks.extend(new_chunks)
                ok_files.append((fname, len(new_chunks)))
            else:
                bad_files.append(fname)

    # Panel de “detalles de procesamiento”
    with st.expander("🔍 View document processing details", expanded=False):
        st.write("### ✅ Loaded files:")
        for fname, n in ok_files:
            st.write(f"✅ {fname} ({n} chunks)")
        st.write(f"\n**Summary:** {len(ok_files)} loaded")
        if not hide_failures and bad_files:
            st.write("### ❌ Skipped/failed files:")
            for fname in bad_files:
                st.write(f"❌ {fname}")

    return chunks


def build_index(chunks: List[Chunk]) -> Tuple[TfidfVectorizer, any]:
    texts = [c.text for c in chunks]
    if not texts:
        vectorizer = TfidfVectorizer(stop_words=None)
        vectorizer.fit(["dummy"])
        matrix = vectorizer.transform(["dummy"])
        return vectorizer, matrix
    vectorizer = TfidfVectorizer(stop_words=None, max_df=0.9)
    matrix = vectorizer.fit_transform(texts)
    return vectorizer, matrix

def get_corpus_hash(root_dir: str) -> str:
    hash_str = ''
    supported_ext = {".pdf", ".docx", ".pptx"}
    for dirpath, _, filenames in os.walk(root_dir):
        for fname in sorted(filenames):
            if os.path.splitext(fname)[1].lower() not in supported_ext:
                continue
            abspath = os.path.join(dirpath, fname)
            mtime = os.path.getmtime(abspath)
            size = os.path.getsize(abspath)
            hash_str += f"{abspath}:{mtime}:{size}\n"
    return hashlib.sha256(hash_str.encode()).hexdigest()

def rank_chunks(query: str, vectorizer: TfidfVectorizer, matrix, chunks: List[Chunk], top_k: int = 25) -> List[Tuple[Chunk, float]]:
    if not query.strip():
        return []
    q_vec = vectorizer.transform([query])
    sims = cosine_similarity(q_vec, matrix)[0]
    idx_scores = sorted(enumerate(sims), key=lambda x: x[1], reverse=True)
    results: List[Tuple[Chunk, float]] = []
    for idx, score in idx_scores[: max(top_k * 2, top_k)]:
        if score <= 0.02:
            continue
        results.append((chunks[idx], float(score)))
        if len(results) >= top_k:
            break
    return results


def extract_relevant_sentences(query: str, texts: List[str], max_sentences: int = 6) -> List[str]:
    query_terms = [t for t in re.split(r"\W+", query.lower()) if len(t) > 2]
    candidates: List[Tuple[str, float]] = []
    for block in texts:
        for sent in split_sentences(block):
            low = sent.lower()
            if not low:
                continue
            tf = sum(low.count(t) for t in query_terms)
            if tf == 0:
                continue
            length_penalty = 1.0 + max(0, (len(sent) - 300) / 300.0)
            score = tf / length_penalty
            candidates.append((sent.strip(), score))
    candidates.sort(key=lambda x: x[1], reverse=True)
    unique: List[str] = []
    seen = set()
    for sent, _ in candidates:
        key = sent[:120]
        if key in seen:
            continue
        seen.add(key)
        unique.append(sent)
        if len(unique) >= max_sentences:
            break
    return unique


def load_name_to_link(project_root: str) -> Dict[str, str]:
    excel_path = os.path.join(project_root, 'Docs & Links', 'Docs & Links - Chatbot P&R Hub.xlsx')
    if not os.path.exists(excel_path):
        return {}
    df = pd.read_excel(excel_path)
    name_col = 'NAME DOCUMENT'
    link_col = 'LINKS'
    if name_col not in df.columns or link_col not in df.columns:
        return {}
    mapping = {}
    for _, row in df.iterrows():
        name = str(row[name_col]).strip()
        link = str(row[link_col]).strip()
        if name and link:
            mapping[name] = link
    return mapping


def compose_answer(query: str, ranked: List[Tuple[Chunk, float]], name_to_link: Dict[str, str] = {}) -> Tuple[str, List[str]]:
    if not ranked:
        msg = (
            "Not found in the available information. "
            "A specific reference (document/page or section) would be needed. "
            "Verify the document name or try other keywords."
        )
        return msg, []
    texts = [c.text for c, _ in ranked]
    sentences = extract_relevant_sentences(query, texts, max_sentences=6)
    if len(sentences) == 0:
        msg = (
            "Not found in the available information. "
            "There are no relevant snippets for the current query."
        )
        return msg, []
    if len(sentences) < 3 and len(ranked) >= 2:
        extra_sentences = []
        for c, _ in ranked:
            sents = split_sentences(c.text)
            for s in sents:
                if len(s.strip()) > 40 and s.strip() not in sentences:
                    extra_sentences.append(s.strip())
                if len(sentences) + len(extra_sentences) >= 3:
                    break
            if len(sentences) + len(extra_sentences) >= 3:
                break
        sentences.extend(extra_sentences[: max(0, 3 - len(sentences))])
    sentences = sentences[:6]
    answer = " ".join(sentences).strip()
    return answer, [f"{name_to_link.get(c.source_name, c.source_name)} — {c.location} — {c.id}" for c, _ in ranked]


def format_sources_lines(ranked: List[Tuple[Chunk, float]] , max_items: int = 10, name_to_link: Dict[str, str] = {}) -> List[str]:
    lines: List[str] = []
    seen = set()
    for c, _ in ranked:
        entry = f"{name_to_link.get(c.source_name, c.source_name)} — {c.location} — {c.id}"
        if entry in seen:
            continue
        seen.add(entry)
        lines.append(entry)
        if len(lines) >= max_items:
            break
    return lines


def call_openai_generate(query: str, ranked: List[Tuple[Chunk, float]], max_sentences: int = 5, custom_system_msg: Optional[str] = None, name_to_link: Dict[str, str] = {}) -> Tuple[Optional[str], Optional[str]]:
    max_ctx = 12
    selected = ranked[:max_ctx]
    context_blocks: List[str] = []
    for c, _ in selected:
        link = name_to_link.get(c.source_name, '')
        context_blocks.append(
            f"[ID: {c.id}]\nLink: {link if link else 'N/A'}\nFile: {c.source_name}\nLocation: {c.location}\nContent: {c.text}"
        )
    context = "\n\n---\n\n" + "\n\n---\n\n".join(context_blocks) if context_blocks else ""

    system_msg = custom_system_msg if custom_system_msg is not None else (
        "You are a RAG assistant. Use ONLY the provided context as your source, "
        "without copying full sentences verbatim from chunks (max 10 consecutive words). "
        "Write clearly and cohesively, interpreting the usage context to adapt the response. "
        "Do not invent or extrapolate beyond the context and preserve acronyms EXACTLY as written. "
        "ALWAYS answer in English, regardless of the user's language. "
        "If there is insufficient evidence, return EXACTLY: 'I cannot find information in the provided chunks to answer this.'"
        "Output instructions:\n"
        "- 3 to 5 sentences, neutral and direct style, no lists.\n"
        "- ALWAYS answer in English and adapt wording to the question's context.\n"
        "- End with the literal 'Sources:' and then, as a list, up to 3 lines each as 'Link — Location — Cited IDs' (use the Link if provided in context, otherwise the File).\n"
        "- If insufficient evidence, return EXACTLY: 'I cannot find information in the provided chunks to answer this.'\n\n"
    )

    user_msg = (
        f"Question: {query}\n\nContext:{context}"
    )

    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg},
    ]
    full_input = json.dumps(messages, ensure_ascii=False)

    try:
        api_key = os.environ.get("OPENAI_API_KEY", "").strip()
        if not api_key:
            return (None, full_input)
        client = OpenAI(api_key=api_key)
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0,
                max_tokens=5000,
            )
            return ((resp.choices[0].message.content or "").strip(), full_input)
        except Exception:
            try:
                resp2 = client.responses.create(
                    model="gpt-4o-mini",
                    input=messages,
                    temperature=0,
                    max_output_tokens=5000,
                )
                if hasattr(resp2, "output") and resp2.output and hasattr(resp2.output[0], "content"):
                    parts = resp2.output[0].content
                    if parts and hasattr(parts[0], "text"):
                        return ((parts[0].text or "").strip(), full_input)
            except Exception:
                return (None, full_input)
    except Exception:
        return (None, full_input)

    return (None, full_input)


def dedupe_preserve_order(items: List[str], limit: int = 5) -> List[str]:
    out: List[str] = []
    seen = set()
    for it in items:
        if it in seen:
            continue
        seen.add(it)
        out.append(it)
        if len(out) >= limit:
            break
    return out

def make_clickable_entry(entry: str) -> str:
    parts = entry.split(" — ", 2)
    if len(parts) < 3:
        return html.escape(entry)
    link, location, cited_id = parts
    if link.startswith("http://") or link.startswith("https://"):
        escaped_link = html.escape(link)
        return f'<a href="{escaped_link}" target="_blank" rel="noopener noreferrer">{escaped_link}</a> — {html.escape(location)} — {html.escape(cited_id)}'
    else:
        return html.escape(entry)

def render_sources_pills(lines: List[str]):
    if not lines:
        st.markdown("<div class='sources-wrap'><span class='source-chip'>not specified</span></div>", unsafe_allow_html=True)
        return
    clickable_lines = [make_clickable_entry(line) for line in lines]
    pills = "".join(f"<span class='source-chip'>{pill}</span>" for pill in clickable_lines)
    st.markdown(f"<div class='sources-wrap'>{pills}</div>", unsafe_allow_html=True)

def render_sources_table(lines: List[str]):
    if not lines:
        st.markdown("<p>No sources available.</p>", unsafe_allow_html=True)
        return

    table_html = '<style>.sources-table{width:auto;border-collapse:collapse;margin-top:0.5rem;}.sources-table th,.sources-table td{border:1px solid #e1efe4;padding:0.2rem;text-align:left;font-size:0.7rem;}.sources-table th{background:#f2f7f3;color:#0f3b1f;}.sources-table .link-cell{max-width:100px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;}.sources-table a{color:#0065BD;text-decoration:none;}.sources-table a:hover{text-decoration:underline;}</style><table class="sources-table"><thead><tr><th>Link</th><th>Location</th><th>ID</th></tr></thead><tbody>'

    for line in lines:
        parts = line.split(" — ", 2)
        if len(parts) < 3:
            continue
        link, location, cited_id = parts
        escaped_location = html.escape(location)
        escaped_id = html.escape(cited_id)
        full_link = name_to_link.get(link, '')
        if full_link.startswith('http'):
            escaped_link = html.escape(full_link)
            link_html = f'<a href="{escaped_link}" target="_blank" rel="noopener noreferrer">{escaped_link}</a>'
        else:
            link_html = html.escape(full_link)
        row_html = f'<tr><td class="link-cell">{link_html}</td><td>{escaped_location}</td><td>{escaped_id}</td></tr>'
        table_html += row_html

    table_html += '</tbody></table>'
    st.markdown(table_html, unsafe_allow_html=True)

def render_app() -> None:
    st.set_page_config(page_title="P&R Hub — RAG (CGIAR)", page_icon="🌿", layout="centered")
    apply_cgiar_theme()

    default_system_prompt = """You are a RAG assistant. Use ONLY the provided context as your source, without copying full sentences verbatim from chunks (max 10 consecutive words). Write clearly and cohesively, interpreting the usage context to adapt the response. Do not invent or extrapolate beyond the context and preserve acronyms EXACTLY as written. ALWAYS answer in English, regardless of the user's language. If there is insufficient evidence, return EXACTLY: 'I cannot find information in the provided chunks to answer this.'Output instructions:\n- 3 to 5 sentences, neutral and direct style, no lists.\n- ALWAYS answer in English and adapt wording to the question's context.\n- End with the literal 'Sources:' and then, as a list, each line as 'File — Location — Cited IDs'.\n- If insufficient evidence, return EXACTLY: 'I cannot find information in the provided chunks to answer this.'\n\n"""

    if 'system_prompt' not in st.session_state:
        st.session_state.system_prompt = default_system_prompt

    # Hero header (solid amber + new title)
    st.markdown("""
        <div class="brand-hero">
            <h1>P&R Hub — Document-grounded RAG assistant</h1>
            <p>Answers strictly based on the documents in this project. Always include the <strong>Sources</strong> section.</p>
        </div>
    """, unsafe_allow_html=True)

    # Sidebar: quick help & example (English, with help icon, and your example)
    with st.sidebar:
        st.markdown("### ❓ Quick help")
        st.write("- Place your **PDF/DOCX/PPTX** in the project folder.")
        st.write("- Ask a specific question.")
        st.write("- Adjust the number of chunks if you need more/less context.")

        st.markdown("---")
        with st.expander("Edit System Prompt"):
            st.session_state.system_prompt = st.text_area("Customize the system prompt for the AI:", value=st.session_state.system_prompt, height=300)

    project_root = os.path.dirname(os.path.abspath(__file__))
    name_to_link = load_name_to_link(project_root)
    chunks_file = os.path.join(project_root, 'chunks.xlsx')
    index_file = os.path.join(project_root, 'index.pkl')
    hash_file = os.path.join(project_root, 'corpus_hash.txt')

    current_hash = get_corpus_hash(project_root)

    load_from_cache = False
    if all(os.path.exists(f) for f in [hash_file, chunks_file, index_file]):
        with open(hash_file, 'r') as f:
            saved_hash = f.read().strip()
        if saved_hash == current_hash:
            load_from_cache = True

    if load_from_cache:
        with st.spinner("Loading from cache..."):
            df = pd.read_excel(chunks_file)
            chunks = [Chunk(**row) for row in df.to_dict(orient='records')]
            with open(index_file, 'rb') as f:
                data = pickle.load(f)
            vectorizer = data['vectorizer']
            matrix = data['matrix']
    else:
        with st.spinner("Loading documents..."):
            chunks = load_corpus(project_root)
        df = pd.DataFrame([asdict(c) for c in chunks])
        df.to_excel(chunks_file, index=False)
        with st.spinner("Building index..."):
            vectorizer, matrix = build_index(chunks)
        with open(index_file, 'wb') as f:
            pickle.dump({'vectorizer': vectorizer, 'matrix': matrix}, f)
        with open(hash_file, 'w') as f:
            f.write(current_hash)

    num_docs = len({c.source_path for c in chunks})
    num_chunks = len(chunks)

    # Metric cards
    cols = st.columns(2, gap="small")

    with cols[0]:
        st.markdown(f"""
        <div class="card" style="display:flex;align-items:center;gap:.65rem">
        <span style="font-size:1.35rem">📄</span>
        <div>
            <div style="font-size:.8rem;color:var(--muted)">Loaded documents</div>
            <div style="font-weight:700">{num_docs}</div>
        </div>
        </div>
        """, unsafe_allow_html=True)

    with cols[1]:
        st.markdown(f"""
        <div class="card" style="display:flex;align-items:center;gap:.65rem">
        <span style="font-size:1.35rem">🧩</span>
        <div>
            <div style="font-size:.8rem;color:var(--muted)">Indexed chunks</div>
            <div style="font-weight:700">{num_chunks}</div>
        </div>
        </div>
        """, unsafe_allow_html=True)

    if num_chunks == 0:
        st.warning(
            "No compatible documents (.pdf, .docx, .pptx) were found in the project. "
            "Add files to the existing folders and reload."
        )

    # Search area as a form
    with st.form("ask_form", clear_on_submit=False):
        query = st.text_input("Type your question:", value="", placeholder="e.g., What is the grievance procedure timeline?")
        top_k = st.slider("Number of chunks to consider", min_value=20, max_value=200, value=100, help="Higher values = more recall, slightly slower.")
        submitted = st.form_submit_button("🔎 Search", use_container_width=True)

    if submitted:
        ranked = rank_chunks(query, vectorizer, matrix, chunks, top_k=top_k)

        # Try OpenAI (if API key present) per your rules
        ai_answer, openai_input = call_openai_generate(query, ranked, max_sentences=5, custom_system_msg=st.session_state.get('system_prompt'), name_to_link=name_to_link)

        # Answer card
        if ai_answer is None or not ai_answer.strip():
            answer, sources_all = compose_answer(query, ranked, name_to_link=name_to_link)

            # Keep answer compact
            unavailable = answer.startswith("I cannot find information in the provided chunks to answer this.") or \
                answer.startswith("Not found in the available information.")
            if not unavailable:
                sents = split_sentences(answer)
                if len(sents) > 6:
                    answer = " ".join(sents[:6]).strip()

            st.markdown(f"<div class='card answer-card'>{html.escape(answer)}</div>", unsafe_allow_html=True)

            # Removed direct sources rendering
        else:
            # For AI answer, only show the body, ignore AI sources
            if 'Sources:' in ai_answer:
                body = ai_answer.split('Sources:', 1)[0].strip()
            else:
                body = ai_answer.strip()

            safe_body = html.escape(body).replace("\n", "<br>")
            st.markdown(f"<div class='card answer-card'>{safe_body}</div>", unsafe_allow_html=True)

            # Removed AI sources rendering

        # Always show TF-IDF sources in table
        st.caption("Top sources (from TF-IDF ranking):")

        data = []
        for c, score in ranked[:5]:
            fuente = c.source_name
            full_link = name_to_link.get(c.source_name, '')
            if not (full_link.startswith('http://') or full_link.startswith('https://')):
                full_link = ''
            data.append({'Fuente': fuente, 'Link': full_link})

        df = pd.DataFrame(data)

        st.markdown("<style> div[data-testid=\"stDataFrame\"] {font-size: 12px;} </style>", unsafe_allow_html=True)

        st.dataframe(
            df,
            width=500,
            hide_index=True,
            column_config={
                "Fuente": st.column_config.TextColumn("Fuente", width="medium"),
                "Link": st.column_config.LinkColumn(
                    "Link 🔗",
                    width=110,
                    display_text=r"^(?:https?://)?([^/]+)",
                    max_chars=20,
                ),
            },
        )

        if len(ranked) > 5:
            st.text("Ver todas")

        # Logging interaction

        log_file = os.path.join(project_root, 'Logs', 'interaction_log.xlsx')

        ranked_chunk_details = ', '.join([f"{c.id}:{score:.3f}" for c, score in ranked])

        log_entry = {
            'timestamp': datetime.now().isoformat(),
            'query': query,
            'top_k': top_k,
            'num_ranked': len(ranked),
            'ai_used': bool(ai_answer and ai_answer.strip()),
            'system_prompt': st.session_state.system_prompt,
            'ranked_chunks': ranked_chunk_details,
            'answer': ai_answer if bool(ai_answer and ai_answer.strip()) else answer,
            'sources': ', '.join(format_sources_lines(ranked, max_items=3, name_to_link=name_to_link)),
            'retrieved_chunks_text': '\n\n---\n\n'.join(f"[ID: {c.id}]\n{c.text}" for c, _ in ranked),
            'openai_input': openai_input or ''
        }

        df_log = pd.DataFrame([log_entry])

        if os.path.exists(log_file):
            existing = pd.read_excel(log_file)
            df_log = pd.concat([existing, df_log], ignore_index=True)

        os.makedirs(os.path.dirname(log_file), exist_ok=True)
        df_log.to_excel(log_file, index=False)

    # Footer
    st.markdown(
        "<div class='app-footer'>Prototype · CGIAR-inspired UI · © 2025</div>",
        unsafe_allow_html=True
    )


if __name__ == "__main__":
    render_app()